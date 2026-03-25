from pathlib import Path
from typing import Dict, List, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from app.mock_assets.generator import MockAssetGenerator
from app.models.schemas import PresentationPlan

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
TITLE_COLOR = RGBColor(22, 37, 66)
ACCENT_COLOR = RGBColor(47, 84, 150)
PAGE_BG = RGBColor(246, 248, 252)
PANEL_BG = RGBColor(255, 255, 255)
BOX_COLOR = RGBColor(234, 240, 248)
BORDER_COLOR = RGBColor(186, 204, 228)
TEXT_COLOR = RGBColor(65, 74, 89)
MUTED_TEXT = RGBColor(114, 123, 140)
LIGHT_LINE = RGBColor(220, 228, 238)


def _add_background(slide, cover: bool = False) -> None:
    background = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    background.fill.solid()
    background.fill.fore_color.rgb = PAGE_BG if not cover else RGBColor(243, 246, 251)
    background.line.color.rgb = background.fill.fore_color.rgb

    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(13.333),
        Inches(0.22 if not cover else 0.4),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT_COLOR
    accent.line.color.rgb = ACCENT_COLOR

    if cover:
        side_panel = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(8.9),
            Inches(0.4),
            Inches(4.433),
            Inches(7.1),
        )
        side_panel.fill.solid()
        side_panel.fill.fore_color.rgb = RGBColor(229, 236, 246)
        side_panel.line.color.rgb = side_panel.fill.fore_color.rgb


def _add_footer(slide, page_number: int, total_pages: int, report_title: str) -> None:
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.5),
        Inches(6.92),
        Inches(12.3),
        Inches(0.03),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = LIGHT_LINE
    line.line.color.rgb = LIGHT_LINE

    left_box = slide.shapes.add_textbox(Inches(0.65), Inches(6.98), Inches(6.4), Inches(0.25))
    left_paragraph = left_box.text_frame.paragraphs[0]
    left_paragraph.text = report_title
    left_paragraph.font.size = Pt(9)
    left_paragraph.font.color.rgb = MUTED_TEXT

    right_box = slide.shapes.add_textbox(Inches(11.5), Inches(6.98), Inches(1.2), Inches(0.25))
    right_paragraph = right_box.text_frame.paragraphs[0]
    right_paragraph.text = "{0:02d} / {1:02d}".format(page_number, total_pages)
    right_paragraph.alignment = PP_ALIGN.RIGHT
    right_paragraph.font.size = Pt(9)
    right_paragraph.font.bold = True
    right_paragraph.font.color.rgb = MUTED_TEXT


def _add_title(slide, title: str) -> None:
    tag_box = slide.shapes.add_textbox(Inches(0.62), Inches(0.42), Inches(1.6), Inches(0.22))
    tag_paragraph = tag_box.text_frame.paragraphs[0]
    tag_paragraph.text = "DFM REPORT"
    tag_paragraph.font.size = Pt(9)
    tag_paragraph.font.bold = True
    tag_paragraph.font.color.rgb = ACCENT_COLOR

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.62), Inches(8.8), Inches(0.46))
    title_paragraph = title_box.text_frame.paragraphs[0]
    title_paragraph.text = title
    title_paragraph.font.size = Pt(24)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = TITLE_COLOR

    divider = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.62),
        Inches(1.1),
        Inches(1.15),
        Inches(0.08),
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = ACCENT_COLOR
    divider.line.color.rgb = ACCENT_COLOR


def _add_panel(slide, left: float, top: float, width: float, height: float, fill_color: RGBColor = PANEL_BG):
    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = fill_color
    panel.line.color.rgb = BORDER_COLOR
    return panel


def _resolve_image_path(asset_generator: MockAssetGenerator, slide_title: str, payload: Dict) -> Optional[str]:
    image_path = payload.get("image_path")
    if image_path and Path(image_path).exists():
        return str(Path(image_path).resolve())

    image_kind = payload.get("image_kind")
    if not image_kind:
        return None

    return asset_generator.ensure_asset(
        kind=image_kind,
        title=payload.get("image_title", slide_title),
        subtitle="Auto-generated placeholder visual",
        labels=payload.get("image_labels", []),
        filename_hint=slide_title,
    )


def _add_image(slide, image_path: Optional[str], left: float, top: float, width: float, height: float, caption: Optional[str] = None) -> None:
    if not image_path:
        return

    panel = _add_panel(slide, left, top, width, height)
    panel.line.color.rgb = BORDER_COLOR
    picture_margin = 0.12
    caption_space = 0.35 if caption else 0
    slide.shapes.add_picture(
        image_path,
        Inches(left + picture_margin),
        Inches(top + picture_margin),
        width=Inches(width - picture_margin * 2),
        height=Inches(height - picture_margin * 2 - caption_space),
    )

    if caption:
        caption_box = slide.shapes.add_textbox(
            Inches(left + 0.15),
            Inches(top + height - 0.3),
            Inches(width - 0.3),
            Inches(0.16),
        )
        paragraph = caption_box.text_frame.paragraphs[0]
        paragraph.text = caption
        paragraph.font.size = Pt(9)
        paragraph.font.color.rgb = MUTED_TEXT


def _add_text_block(
    slide,
    lines: List[str],
    left: float,
    top: float,
    width: float,
    height: float,
    font_size: int = 16,
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0

    for index, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = u"\u2022 " + line
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = TEXT_COLOR
        paragraph.space_after = Pt(7)


def _estimate_column_widths(headers: List[str], rows: List[List[str]], total_width: float) -> List[float]:
    scores = []
    for col_index, header in enumerate(headers):
        max_len = len(str(header))
        for row in rows:
            if col_index < len(row):
                max_len = max(max_len, len(str(row[col_index])))
        scores.append(max(8, min(max_len, 36)))

    total_score = float(sum(scores)) or 1.0
    widths = [max(0.95, total_width * score / total_score) for score in scores]
    scale = total_width / sum(widths)
    return [width * scale for width in widths]


def _add_table(
    slide,
    headers: List[str],
    rows: List[List[str]],
    left: float = 0.6,
    top: float = 1.45,
    width: float = 12.1,
    height: float = 5.15,
) -> None:
    table_panel = _add_panel(slide, left, top, width, height)
    table_panel.fill.fore_color.rgb = PANEL_BG

    row_count = max(2, len(rows) + 1)
    col_count = len(headers)
    table = slide.shapes.add_table(
        row_count,
        col_count,
        Inches(left + 0.08),
        Inches(top + 0.08),
        Inches(width - 0.16),
        Inches(height - 0.16),
    ).table

    column_widths = _estimate_column_widths(headers, rows, width - 0.16)
    for col_index, col_width in enumerate(column_widths):
        table.columns[col_index].width = Inches(col_width)

    for col_index, header in enumerate(headers):
        cell = table.cell(0, col_index)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT_COLOR
        cell.margin_left = 60000
        cell.margin_right = 60000
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        paragraph.font.size = Pt(11)

    for row_index in range(1, row_count):
        row_values = rows[row_index - 1] if row_index - 1 < len(rows) else [""] * col_count
        for col_index in range(col_count):
            value = row_values[col_index] if col_index < len(row_values) else ""
            cell = table.cell(row_index, col_index)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(248, 250, 253) if row_index % 2 else RGBColor(241, 245, 250)
            cell.margin_left = 50000
            cell.margin_right = 50000
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(10.5)
            paragraph.font.color.rgb = TEXT_COLOR
            paragraph.alignment = PP_ALIGN.LEFT
            cell.text_frame.word_wrap = True


def _add_stat_cards(slide, stats: List[str]) -> None:
    for index, stat in enumerate(stats):
        left = 0.8 + index * 2.42
        card = _add_panel(slide, left, 5.75, 2.18, 0.86, fill_color=RGBColor(235, 241, 249))
        card.line.color.rgb = RGBColor(196, 210, 230)
        text_frame = card.text_frame
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        paragraph = text_frame.paragraphs[0]
        paragraph.text = stat
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.size = Pt(14)
        paragraph.font.bold = True
        paragraph.font.color.rgb = TITLE_COLOR


def _render_cover(slide, title: str, payload: Dict, image_path: Optional[str]) -> None:
    label_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.9), Inches(2.2), Inches(0.25))
    label_paragraph = label_box.text_frame.paragraphs[0]
    label_paragraph.text = "DESIGN FOR MANUFACTURING"
    label_paragraph.font.size = Pt(10)
    label_paragraph.font.bold = True
    label_paragraph.font.color.rgb = ACCENT_COLOR

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.25), Inches(6.6), Inches(1.2))
    title_paragraph = title_box.text_frame.paragraphs[0]
    title_paragraph.text = title
    title_paragraph.font.size = Pt(28)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = TITLE_COLOR

    subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(6.6), Inches(0.4))
    subtitle_paragraph = subtitle_box.text_frame.paragraphs[0]
    subtitle_paragraph.text = payload["subtitle"]
    subtitle_paragraph.font.size = Pt(15)
    subtitle_paragraph.font.color.rgb = MUTED_TEXT

    meta_panel = _add_panel(slide, 0.8, 2.9, 6.55, 2.25, fill_color=PANEL_BG)
    meta_frame = meta_panel.text_frame
    meta_frame.clear()
    for index, meta_row in enumerate(payload.get("report_meta", [])):
        paragraph = meta_frame.paragraphs[0] if index == 0 else meta_frame.add_paragraph()
        paragraph.text = "{0}: {1}".format(meta_row[0], meta_row[1])
        paragraph.font.size = Pt(13)
        paragraph.font.color.rgb = TEXT_COLOR
        paragraph.space_after = Pt(6)

    _add_stat_cards(slide, payload["stats"])
    _add_image(slide, image_path, 9.25, 1.0, 3.45, 4.95, caption="Overview placeholder image")


def _render_bullets(slide, payload: Dict, image_path: Optional[str]) -> None:
    if image_path:
        _add_panel(slide, 0.6, 1.35, 7.2, 5.25)
        _add_text_block(slide, payload["bullets"], left=0.85, top=1.68, width=6.6, height=4.7)
    else:
        _add_panel(slide, 0.6, 1.35, 12.1, 5.25)
        _add_text_block(slide, payload["bullets"], left=0.9, top=1.7, width=11.4, height=4.7)
    if image_path:
        _add_image(slide, image_path, 8.05, 1.35, 4.65, 5.25, caption="Mock supporting visual")


def _render_process(slide, payload: Dict, image_path: Optional[str]) -> None:
    process_panel = _add_panel(slide, 0.6, 1.35, 7.45, 5.25)
    steps = payload["steps"]
    for index, step in enumerate(steps):
        card = _add_panel(slide, 0.9, 1.8 + index * 0.82, 5.55, 0.58, fill_color=RGBColor(240, 245, 251))
        paragraph = card.text_frame.paragraphs[0]
        paragraph.text = "{0}. {1}".format(index + 1, step)
        paragraph.font.size = Pt(15)
        paragraph.font.bold = True
        paragraph.font.color.rgb = TITLE_COLOR

        if index < len(steps) - 1:
            arrow = slide.shapes.add_textbox(Inches(6.55), Inches(2.0 + index * 0.82), Inches(0.35), Inches(0.2))
            arrow_paragraph = arrow.text_frame.paragraphs[0]
            arrow_paragraph.text = "v"
            arrow_paragraph.font.size = Pt(18)
            arrow_paragraph.font.bold = True
            arrow_paragraph.font.color.rgb = ACCENT_COLOR
            arrow_paragraph.alignment = PP_ALIGN.CENTER

    _add_image(slide, image_path, 8.3, 1.35, 4.4, 5.25, caption="Workflow placeholder image")


def _render_structure(slide, payload: Dict, image_path: Optional[str]) -> None:
    structure_panel = _add_panel(slide, 0.6, 1.35, 7.55, 5.25)

    root_box = _add_panel(slide, 2.5, 1.8, 3.6, 0.72, fill_color=BOX_COLOR)
    root_frame = root_box.text_frame
    root_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    root_paragraph = root_frame.paragraphs[0]
    root_paragraph.text = payload["root"]
    root_paragraph.font.size = Pt(18)
    root_paragraph.font.bold = True
    root_paragraph.font.color.rgb = TITLE_COLOR
    root_paragraph.alignment = PP_ALIGN.CENTER

    positions = [(0.9, 3.4), (2.35, 3.4), (3.8, 3.4), (5.25, 3.4), (6.7, 3.4)]
    for (left, top), child in zip(positions, payload["children"][:5]):
        connector = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(left + 0.52),
            Inches(2.65),
            Inches(0.03),
            Inches(0.48),
        )
        connector.fill.solid()
        connector.fill.fore_color.rgb = ACCENT_COLOR
        connector.line.color.rgb = ACCENT_COLOR

        child_box = _add_panel(slide, left, top, 1.15, 0.95, fill_color=RGBColor(246, 249, 253))
        child_frame = child_box.text_frame
        child_frame.word_wrap = True
        child_paragraph = child_frame.paragraphs[0]
        child_paragraph.text = child
        child_paragraph.font.size = Pt(11)
        child_paragraph.font.bold = True
        child_paragraph.font.color.rgb = TEXT_COLOR
        child_paragraph.alignment = PP_ALIGN.CENTER

    _add_image(slide, image_path, 8.4, 1.35, 4.3, 5.25, caption="Overview placeholder image")


def _render_module_overview(slide, payload: Dict, image_path: Optional[str]) -> None:
    modules = payload["modules"][:4]
    positions = [(0.6, 1.35), (6.55, 1.35), (0.6, 3.2), (6.55, 3.2)]
    for (left, top), module in zip(positions, modules):
        card = _add_panel(slide, left, top, 5.75, 1.45, fill_color=RGBColor(245, 248, 252))
        frame = card.text_frame
        title_paragraph = frame.paragraphs[0]
        title_paragraph.text = "{0} ({1})".format(module["name"], module["count"])
        title_paragraph.font.size = Pt(17)
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = TITLE_COLOR
        body_paragraph = frame.add_paragraph()
        body_paragraph.text = module["preview"]
        body_paragraph.font.size = Pt(12)
        body_paragraph.font.color.rgb = TEXT_COLOR
        body_paragraph.space_before = Pt(6)

    if image_path:
        _add_image(slide, image_path, 2.3, 5.0, 8.7, 1.45, caption="Module overview placeholder image")


def _render_module_detail(slide, payload: Dict, image_path: Optional[str]) -> None:
    summary_panel = _add_panel(slide, 0.6, 1.35, 3.35, 2.55, fill_color=RGBColor(245, 248, 252))
    summary_frame = summary_panel.text_frame
    summary_title = summary_frame.paragraphs[0]
    summary_title.text = "Summary"
    summary_title.font.size = Pt(17)
    summary_title.font.bold = True
    summary_title.font.color.rgb = TITLE_COLOR
    for bullet in payload["summary"]:
        paragraph = summary_frame.add_paragraph()
        paragraph.text = u"\u2022 " + bullet
        paragraph.font.size = Pt(11.5)
        paragraph.font.color.rgb = TEXT_COLOR
        paragraph.space_before = Pt(4)

    _add_image(slide, image_path, 4.15, 1.35, 8.55, 2.55, caption="Module placeholder image")
    _add_table(slide, payload["table_rows"][0], payload["table_rows"][1:], left=0.6, top=4.15, width=12.1, height=2.4)


def _render_architecture(slide, payload: Dict, image_path: Optional[str]) -> None:
    panel = _add_panel(slide, 0.6, 1.35, 6.55, 5.25)
    top = 1.75
    for layer in payload["layers"]:
        layer_box = _add_panel(slide, 0.95, top, 5.85, 0.72, fill_color=RGBColor(241, 246, 251))
        paragraph = layer_box.text_frame.paragraphs[0]
        paragraph.text = layer
        paragraph.font.size = Pt(15)
        paragraph.font.bold = True
        paragraph.font.color.rgb = TITLE_COLOR
        paragraph.alignment = PP_ALIGN.CENTER
        if top < 4.7:
            connector = slide.shapes.add_textbox(Inches(3.78), Inches(top + 0.75), Inches(0.2), Inches(0.18))
            connector_paragraph = connector.text_frame.paragraphs[0]
            connector_paragraph.text = "v"
            connector_paragraph.font.size = Pt(18)
            connector_paragraph.font.bold = True
            connector_paragraph.font.color.rgb = ACCENT_COLOR
            connector_paragraph.alignment = PP_ALIGN.CENTER
        top += 1.0

    _add_image(slide, image_path, 7.45, 1.35, 5.25, 5.25, caption="Software architecture placeholder image")


def build_presentation(plan: PresentationPlan, output_path: str) -> str:
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    output_file = Path(output_path)
    output_file.parent.mkdir(parents=True, exist_ok=True)
    asset_dir = output_file.parent / "assets" / output_file.stem
    asset_generator = MockAssetGenerator(asset_dir)
    total_pages = len(plan.slides)

    for page_number, slide_spec in enumerate(plan.slides, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        image_path = _resolve_image_path(asset_generator, slide_spec.title, slide_spec.payload)

        if slide_spec.slide_type == "cover":
            _add_background(slide, cover=True)
            _render_cover(slide, slide_spec.title, slide_spec.payload, image_path)
        else:
            _add_background(slide)
            _add_title(slide, slide_spec.title)

            if slide_spec.slide_type == "bullets":
                _render_bullets(slide, slide_spec.payload, image_path)
            elif slide_spec.slide_type == "process":
                _render_process(slide, slide_spec.payload, image_path)
            elif slide_spec.slide_type == "structure":
                _render_structure(slide, slide_spec.payload, image_path)
            elif slide_spec.slide_type == "table":
                _add_table(slide, slide_spec.payload["headers"], slide_spec.payload["rows"])
            elif slide_spec.slide_type == "module_overview":
                _render_module_overview(slide, slide_spec.payload, image_path)
            elif slide_spec.slide_type == "module_detail":
                _render_module_detail(slide, slide_spec.payload, image_path)
            elif slide_spec.slide_type == "architecture":
                _render_architecture(slide, slide_spec.payload, image_path)

        _add_footer(slide, page_number=page_number, total_pages=total_pages, report_title=plan.report_title)

    prs.save(output_file)
    return str(output_file.resolve())
